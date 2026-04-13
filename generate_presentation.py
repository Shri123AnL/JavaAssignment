from pptx import Presentation
from pptx.util import Inches

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide_1.shapes.title
    subtitle = slide_1.placeholders[1]
    title.text = "Algorithm Analysis and Design"
    subtitle.text = "A Professional Overview"

    # Slide 2: Introduction to Algorithms
    slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide_2.shapes.title
    content = slide_2.placeholders[1]
    title.text = "Introduction"
    content.text = "An algorithm is a step-by-step procedure for calculations and data processing."

    # Slide 3: Algorithm Complexity
    slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide_3.shapes.title
    content = slide_3.placeholders[1]
    title.text = "Algorithm Complexity"
    content.text = "Understanding time and space complexity helps in analyzing the efficiency of algorithms."

    # Slide 4: Types of Algorithms
    slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide_4.shapes.title
    content = slide_4.placeholders[1]
    title.text = "Types of Algorithms"
    content.text = "1. Divide and Conquer\n2. Greedy Algorithms\n3. Dynamic Programming"

    # Slide 5: Conclusion and Applications
    slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide_5.shapes.title
    content = slide_5.placeholders[1]
    title.text = "Conclusion"
    content.text = "Algorithms are essential in computer science for effective problem-solving."

    # Save the presentation
    prs.save('Algorithm_Analysis_and_Design_Presentation.pptx')

if __name__ == "__main__":
    create_presentation()